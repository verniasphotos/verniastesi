from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a presentation object
prs = Presentation()

# Standard Slide Layouts
# 0: Title Slide
# 1: Title and Content
# 5: Title Only
# 6: Blank

def add_title_slide(prs, title_text, subtitle_text, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle.text = subtitle_text
    
    if os.path.exists(image_path):
        # Add small logo at the top or bottom
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(0.5), height=Inches(1.5))
    return slide

def add_content_slide(prs, title_text, content_text, image_path=None, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    
    # Check if there is text content
    if content_text:
        body.text = content_text
        for p in body.text_frame.paragraphs:
            p.font.size = Pt(20)
    else:
        # If no text, we just use the space for the image
        sp = body.element
        sp.getparent().remove(sp)
    
    # Add Image
    if image_path and os.path.exists(image_path):
        if content_text:
            # Place image on the right
            left = Inches(4.5)
            top = Inches(1.8)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, height=height)
            # Resize text box to be smaller
            body.width = Inches(4.0)
        else:
            # Center image
            left = Inches(1.5)
            top = Inches(1.5)
            height = Inches(5.0)
            slide.shapes.add_picture(image_path, left, top, height=height)
            
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes
        
    return slide

def add_large_image_slide(prs, title_text, image_path, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    title = slide.shapes.title
    title.text = title_text
    
    if image_path and os.path.exists(image_path):
        # Best effort center
        try:
            pic = slide.shapes.add_picture(image_path, Inches(1.0), Inches(1.5), height=Inches(5.5))
            # Center horizontally roughly
            pic.left = int((prs.slide_width - pic.width) / 2)
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
            
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes
        
    return slide

# Base Directory
base_dir = "/Users/vernias/Desktop/verniastesi"

# --- Slide 1: Titolo ---
slide1 = add_title_slide(prs, 
                         "Progettazione e verifica di un'architettura 6G per il tracking di droni indoor assistita da RIS",
                         "Candidato: Matteo Verniani\nRelatore: Chiar.mo Prof. [Nome]",
                         os.path.join(base_dir, "Report/Images/logo-unifi.png"))
# Manually add notes to slide 1
notes_slide = slide1.notes_slide
notes_slide.notes_text_frame.text = "Buongiorno alla commissione. Sono Matteo Verniani e oggi presento il mio lavoro di tesi incentrato sulla progettazione e validazione di un’architettura 6G avanzata per il tracking di droni autonomi in ambienti indoor complessi, sfruttando la tecnologia delle Reconfigurable Intelligent Surfaces o RIS. Ringrazio il mio relatore per il supporto."

# --- Slide 2: Contesto ---
add_content_slide(prs, 
                  "Contesto: Logistica 4.0 e Navigazione VNA", 
                  "• Applicazioni in magazzini intensivi VNA (Very Narrow Aisle)\n• Uso massiccio di droni per material handling e inventario\n• Assenza di segnale GPS (Indoor)\n• Severi ostacoli metallici (scaffalature)",
                  os.path.join(base_dir, "warehouse_drone_problem.png"),
                  notes="Il contesto applicativo è la Logistica 4.0, in scenari come i magazzini intensivi VNA, dove si fa largo uso di droni. In questi ambienti indoor e metallici, il GPS è inefficace. Inoltre, la presenza massiccia di scaffalature crea severi ostacoli fisici.")

# --- Slide 3: NLoS ---
add_content_slide(prs,
                  "Il Problema: Non-Line-of-Sight (NLoS)",
                  "• Ostruzione profonda del segnale radio\n• Perdita del link RF diretto con la Base Station\n• Errore di posizionamento divergente\n• Impossibilità di navigazione autonoma sicura",
                  os.path.join(base_dir, "nlos_vna_concept.png"),
                  notes="Il problema fondamentale è l'ostruzione della Line-of-Sight. Quando il drone si sposta dietro una scaffalatura, entriamo in condizioni di Non-Line-of-Sight profondo. I normali sistemi radio perdono il segnale, l'errore aumenta e la navigazione autonoma diventa impossibile.")

# --- Slide 4: 6G e RIS ---
add_content_slide(prs,
                  "La Soluzione 6G: Reconfigurable Intelligent Surfaces",
                  "• 6G: Paradigma di comunicazione intelligente\n• RIS: Superfici programmabili per reindirizzare le onde radio\n• Creazione di una 'Virtual Line-of-Sight' (V-LoS)\n• Ripristino della visibilità radio oltre l'ostacolo",
                  os.path.join(base_dir, "active_ris.png"),
                  notes="La soluzione risiede nel paradigma emergente del 6G: l'adozione delle RIS. Queste superfici intelligenti permettono di reindirizzare le onde radio, aggirando fisicamente l'ostacolo. Si crea così una 'Virtual Line-of-Sight', riportando il drone nel range di tracking.")

# --- Slide 5: Architettura ---
add_large_image_slide(prs,
                      "L'Architettura Globale (SDN / O-RAN)",
                      os.path.join(base_dir, "Architettura_Blocchi_Semplificata.png"),
                      notes="Il cuore di questa tesi è la progettazione di un simulatore Digital Twin basato su paradigma SDN (Software-Defined Networking). Un Server Centrale riceve la telemetria, esegue gli algoritmi di tracking in real-time e orchestra le RIS presenti nell'impianto per garantire sempre la miglior connessione.")

# --- Slide 6: Tracking Ibrido ---
# Assuming you have a schematic for this, using an empty/generic layout if no specific image
add_content_slide(prs,
                  "Motore di Tracking: Fusione Ibrida EKF-LSTM",
                  "• Extended Kalman Filter (EKF) per le traiettorie lineari e aggregazione dati\n• Limiti: EKF 'linearizza' in zone NLoS sfasando sulle curve\n• Long Short-Term Memory (LSTM) per stima puramente cinematica in assenza di segnale\n• Predizione Non-Lineare su pattern di volo",
                  None,
                  notes="L'intelligenza di posizionamento è ibrida. In LoS, utilizziamo l'EKF per fondere inerziale e radio. Tuttavia, in NLoS prolungato, il modello passa il controllo a una LSTM, addestrata sulle traiettorie del layout, per una stima puramente cinematica e non lineare.")

# --- Slide 7: Make-Before-Break ---
add_content_slide(prs,
                  "Controllo Predittivo: Make-Before-Break",
                  "• Combinazione di SDN e pre-triggering LSTM\n• Attivazione della RIS preventiva\n• Configurazione ottimale della fase prima dell'ingresso in zona buia (NLoS)\n• Tracking 'Seamless' (senza interruzioni)",
                  None,
                  notes="Questa architettura abilita un controllo SDN predittivo: 'Make-Before-Break'. Invece di reagire a una disconnessione, si prevede il movimento del drone. Il controller accende la RIS ottima msec prima che il drone entri nella zona NLoS, garantendo tracking continuo.")

# --- Slide 8: Setup ---
add_content_slide(prs,
                  "Setup Sperimentale: Il Simulatore 6G",
                  "• Sviluppo interamente in Python (motore DDES)\n• Calcolo matriciale avanzato tramite Numba\n• Comunicazione distribuita con protocollo gRPC\n• Modellazione realistica canale InF-DH (Indoor Factory) ai 5.9GHz",
                  None,
                  notes="Per la validazione, ho sviluppato un simulatore ad eventi implementato in Python. Esegue il calcolo vettoriale pesante con Numba e comunica mediante gRPC, modellando con estremo realismo un canale InF-DH ai 5.9GHz. La validazione si struttura su 4 test.")

# --- Slide 9: Test 1 ---
add_large_image_slide(prs,
                      "Test 1 (Baseline): Fallimento Coasting EKF in NLoS",
                      os.path.join(base_dir, "Test_1_RMSE_Temporale.png"),
                      notes="Test 1: il comportamento Baseline senza RIS. Quando il drone entra in corridoio con ostruzione totale, il Filtro EKF va in deriva 'Coasting', e l'errore di posizionamento sale esponenzialmente portando inevitabilmente a un errore catastrofico o collisione.")

# --- Slide 10: Limite Lineare vs Deep Learning ---
add_large_image_slide(prs,
                      "Confronto Deriva Lineare EKF vs Predizione LSTM",
                      os.path.join(base_dir, "Test_4.2_Traiettoria_EKF_vs_LSTM.png"),
                      notes="Il fallimento in navigazione cieca avviene perché il predittore inerziale EKF procede linearmente. La nostra LSTM subentra qui. Osserviamo come l'EKF linearizzi ciecamente l'ostacolo (linea rossa), mentre il Deep Learning replichi aderentemente la virata nel corridoio (blu vs verde).")

# --- Slide 11: Test 2 ---
add_large_image_slide(prs,
                      "Test 2: Successo del Tracciamento Assistito da RIS attive",
                      os.path.join(base_dir, "Test_2_RMSE_Temporale.png"),
                      notes="Test 2: accendiamo l'architettura 6G. Al presentarsi del NLoS, SDN attiva preventivamente le RIS. L'errore (RMSE) rientra e rimane saldamente sotto il metro di tolleranza pattuito, mantenendo stabile il volo e completando la missione in sicurezza.")

# --- Slide 12: Test 3 ---
add_large_image_slide(prs,
                      "Test 3: Stress-Test della Latenza di Rete",
                      os.path.join(base_dir, "Test_3_Stress_Latenza.png"),
                      notes="L'architettura è sensibile ai ritardi di rete. Abbiamo stressato il sistema introducendo latenza nel loop di controllo. Il limite fisico impone una soglia restrittiva inferiore a ~50ms per evitare il Beam Misalignment (mancare il drone col raggio della RIS).")

# --- Slide 13: Sfida Energetica ---
add_content_slide(prs,
                  "Oltre le Prestazioni: La Sfida dell'Efficienza (Green 6G)",
                  "• Sviluppo di infrastruttura solida (Aggancio e Tracking riusciti)\n• Problema nascente: consumo energetico delle stazioni e RIS attive al 100%\n• Costi e dissipazioni termiche incompatibili con il paradigma IoT e industriale\n• Necessità di soluzioni termodinamicamente equilibrate",
                  None,
                  notes="Mantenere un aggancio risolve la precisione, ma le stazioni e le RIS al massimo della potenza comportano costi/consumi insostenibili. La sfida finale del lavoro consiste nel rendere questo ecosistema sostenibile, ottimizzando il risparmio energetico: il Green 6G.")

# --- Slide 14: Test 4 ---
add_large_image_slide(prs,
                      "Test 4: Ottimizzazione Green 6G tramite Dinkelbach",
                      os.path.join(base_dir, "Test_4.6_AEE_vs_Ps_Dinkelbach.png"),
                      notes="Ho usato l'algoritmo di Dinkelbach e la matrice BD-RIS per l'Efficienza Energetica Assoluta (AEE). Il network ricalcola il punto ottimo termodinamico ad ogni step. Troviamo così il 'Pareto-Front' evidenziato: picco di massima efficienza radiativa con potenze moderate.")

# --- Slide 15: Conclusioni ---
add_content_slide(prs,
                  "Conclusioni e Goal Raggiunti",
                  "• Ripristino di tracking sub-metrico in ambiente industriale blindato (NLoS)\n• Strategia Make-Before-Break abilitata dalla predizione LSTM su rete SDN\n• Dimostrazione dell'efficacia delle BD-RIS (Beyond Diagonal) attive\n• Bilanciamento Pareto-ottimale raggiunto: performance 6G e sostenibilità Green",
                  None,
                  notes="In conclusione, l'architettura progettata ha permesso di riacquisire il tracking indoor NLoS, sfruttando un approccio SDN predittivo assistito da LSTM ('Make-Before-Break'). Infine, l'algoritmo di Dinkelbach su BD-RIS consente di coniugare le elevate prestazioni 6G con l'efficienza Green, centrale per le future reti. Grazie.")


prs.save(os.path.join(base_dir, "Presentazione_Tesi.pptx"))
print("Presentazione_Tesi.pptx creata con successo!")
